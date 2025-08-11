"""
pptx_inconsistency_detector.py

A self-contained Python CLI tool to analyze a PowerPoint (.pptx) and flag factual or logical
inconsistencies across slides.

Features:
- Extracts text from slides (python-pptx) and runs OCR on embedded images (pytesseract + Pillow).
- Parses numeric facts (currency, percent, plain numbers) and dates from slide text.
- Normalizes numeric values for comparison.
- Builds candidate pairs of claims and compares them using a) a rule-based comparator and b) an
  external LLM comparator (Gemini 2.5 Flash) if API credentials and endpoint are provided.
- Produces a JSON report with flagged inconsistencies and a human-readable summary printed to stdout.

Limitations:
- OCR accuracy depends on image quality and Tesseract configuration.
- Semantic contradiction detection depends on the LLM (if used). Without LLM, only numeric/date
  checks and simple textual contradictions are detected.
- This is a prototype. For production you'd add more advanced NER, coreference resolution, and
  more robust batching/rate-limiting for API calls.

Dependencies:
- python-pptx
- pillow
- pytesseract
- requests
- dateparser
- regex (the builtin 're' is used but 'regex' can be swapped in)

Install via:
    pip install python-pptx pillow pytesseract requests dateparser

Usage example:
    python pptx_inconsistency_detector.py --input deck.pptx --out report.json

To enable Gemini comparator, set environment variables:
    GEMINI_API_KEY, GEMINI_ENDPOINT

GEMINI_ENDPOINT should be the full REST API endpoint you will POST to. The exact Google
GenAI endpoint/SDK varies; this script expects a JSON response from the model (or text which
it will attempt to parse). If you don't have a Gemini endpoint, the script will still run
rule-based checks.

"""

import os
import sys
import json
import re
import argparse
import math
import time
import tempfile
from collections import defaultdict
from dataclasses import dataclass, asdict
from typing import List, Dict, Any, Optional, Tuple

from pptx import Presentation
from PIL import Image
import pytesseract
import requests
import dateparser
from dotenv import load_dotenv
load_dotenv()

# --------------------------- Data classes ---------------------------

@dataclass
class SlideContent:
    slide_no: int
    raw_text: str
    image_paths: List[str]

@dataclass
class Fact:
    slide_no: int
    kind: str  # 'number', 'percent', 'currency', 'date', 'claim'
    raw: str
    normalized: Any
    context: str

@dataclass
class Issue:
    type: str
    slide_a: int
    slide_b: int
    snippet_a: str
    snippet_b: str
    details: Dict[str, Any]
    confidence: float
    explanation: str

# --------------------------- Extraction ---------------------------

def extract_from_pptx(path: str, extract_images_to: Optional[str] = None) -> List[SlideContent]:
    """Extract text and images from a pptx file.

    Returns list of SlideContent objects.
    """
    prs = Presentation(path)
    slides: List[SlideContent] = []
    tmpdir = extract_images_to or tempfile.mkdtemp(prefix="pptx_images_")

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        image_paths = []

        # Extract text from shapes
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text:
                    texts.append(shape.text.strip())
            except Exception:
                # not all shapes support has_text_frame
                pass

            # Save picture shapes
            if shape.shape_type == 13:  # 13 -> picture
                try:
                    image = shape.image
                    img_bytes = image.blob
                    ext = image.ext
                    img_filename = os.path.join(tmpdir, f"slide_{i}_img_{len(image_paths)}.{ext}")
                    with open(img_filename, "wb") as f:
                        f.write(img_bytes)
                    image_paths.append(img_filename)
                except Exception:
                    pass

        # Some slides may be entirely images exported as background shapes; python-pptx can't always
        # access those. As fallback, save slide as image by rendering — python-pptx doesn't render.
        # So we rely on any picture shapes present. For complex cases, user should provide slide images.

        slides.append(SlideContent(slide_no=i, raw_text="\n".join(texts), image_paths=image_paths))

    return slides

# --------------------------- OCR ---------------------------

def ocr_image(image_path: str, lang: str = 'eng') -> str:
    """Run pytesseract OCR on an image and return extracted text."""
    try:
        txt = pytesseract.image_to_string(Image.open(image_path), lang=lang)
        return txt.strip()
    except Exception as e:
        print(f"[WARN] OCR failed on {image_path}: {e}")
        return ""

# --------------------------- Parsing facts ---------------------------

# Regex patterns
CURRENCY_RE = re.compile(r"(?P<prefix>[$£€₹])?\s*(?P<val>[0-9][0-9,\.]*)(?:\s*(?P<scale>[kKmMbB]))?\s*(?P<suffix>USD|EUR|GBP|INR)?", re.IGNORECASE)
PERCENT_RE = re.compile(r"(?P<val>[0-9]{1,3}(?:[\.,][0-9]+)?)\s*%")
NUMBER_RE = re.compile(r"(?<!\w)(?P<val>\d{1,3}(?:[,\d]{3})*(?:\.\d+)?)(?:\s*(?P<scale>k|m|b))?(?!\w)", re.IGNORECASE)
DATE_RES = [
    # Try to match common explicit date mentions
    re.compile(r"(Q[1-4]\s*\d{2,4})", re.IGNORECASE),
    re.compile(r"(\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{2,4})", re.IGNORECASE),
    re.compile(r"(\d{4})"),
]

SIMPLE_CLAIM_RE = re.compile(r"([A-Z][^\.\n]{10,200})")  # naive: sentences starting with capital letter


def normalize_number(value_str: str, scale_hint: Optional[str] = None) -> Optional[float]:
    """Normalize a numeric string into a float. Handles commas and simple scales (k,m,b).

    Returns None if parsing fails.
    """
    if not value_str:
        return None
    s = value_str.replace(',', '').strip()
    try:
        val = float(s)
    except Exception:
        # handle cases like '1.5M'
        m = re.match(r"([0-9]*\.?[0-9]+)\s*([kKmMbB])", value_str)
        if m:
            base = float(m.group(1))
            scale = m.group(2).lower()
            if scale == 'k':
                return base * 1_000
            if scale == 'm':
                return base * 1_000_000
            if scale == 'b':
                return base * 1_000_000_000
        return None

    # apply scale_hint if present
    if scale_hint:
        sh = scale_hint.lower()
        if sh == 'k':
            val *= 1_000
        elif sh == 'm':
            val *= 1_000_000
        elif sh == 'b':
            val *= 1_000_000_000
    return val


def parse_facts_from_text(text: str, slide_no: int) -> List[Fact]:
    facts: List[Fact] = []
    if not text:
        return facts

    # Currency
    for m in CURRENCY_RE.finditer(text):
        val = m.group('val')
        scale = m.group('scale')
        cur = m.group('prefix') or m.group('suffix')
        norm = normalize_number(val, scale_hint=scale)
        kind = 'currency'
        raw = m.group(0)
        facts.append(Fact(slide_no=slide_no, kind=kind, raw=raw, normalized={'value': norm, 'currency': cur}, context=text[:300]))

    # Percent
    for m in PERCENT_RE.finditer(text):
        val = m.group('val').replace(',', '.')
        try:
            norm = float(val)
        except Exception:
            norm = None
        facts.append(Fact(slide_no=slide_no, kind='percent', raw=m.group(0), normalized=norm, context=text[:300]))

    # Plain numbers
    for m in NUMBER_RE.finditer(text):
        # Avoid duplicating currency/percent matches by simple check
        snippet = m.group(0)
        if '%' in snippet:
            continue
        # skip if it's clearly part of a currency
        if CURRENCY_RE.search(snippet):
            continue
        val = m.group('val')
        scale = m.group('scale')
        norm = normalize_number(val, scale_hint=scale)
        facts.append(Fact(slide_no=slide_no, kind='number', raw=snippet, normalized=norm, context=text[:300]))

    # Dates (naive: try dateparser on contexts)
    # We'll scan sentences and try to parse any date-like substring
    sentences = re.split(r'[\n\.]', text)
    for s in sentences:
        s = s.strip()
        if not s:
            continue
        # Try to parse explicit date-like tokens
        dt = dateparser.parse(s, settings={'PREFER_DATES_FROM': 'past'})
        if dt:
            facts.append(Fact(slide_no=slide_no, kind='date', raw=s, normalized=dt.isoformat(), context=s))

    # Simple claim extraction as fallback (non-fact sentences)
    for m in SIMPLE_CLAIM_RE.finditer(text):
        claim = m.group(0).strip()
        if len(claim) > 20:
            facts.append(Fact(slide_no=slide_no, kind='claim', raw=claim, normalized=None, context=claim))

    # Deduplicate by raw text
    unique = []
    seen = set()
    for f in facts:
        if f.raw in seen:
            continue
        seen.add(f.raw)
        unique.append(f)
    return unique

# --------------------------- Pairing & Candidate generation ---------------------------

def group_facts_by_metric(facts: List[Fact]) -> Dict[str, List[Fact]]:
    """Group facts by heuristic metric key. For numeric/currency, group by nearby words.

    This is intentionally heuristic: we use the first few words of context near the match as a key.
    """
    groups: Dict[str, List[Fact]] = defaultdict(list)
    for f in facts:
        key = f.kind
        # for numbers/currency, try to find metric words nearby (e.g., revenue, ARR, users)
        context = f.context or ''
        # simple heuristic: extract the word before and after the raw token in context if available
        try:
            idx = context.find(f.raw)
            if idx >= 0:
                snippet = context[max(0, idx-40): idx+len(f.raw)+40]
            else:
                snippet = context[:80]
        except Exception:
            snippet = context[:80]

        # metric label: first noun-like token in snippet (naive)
        words = re.findall(r"[A-Za-z%$₹£€]+", snippet)
        metric = None
        for w in words[:10]:
            lw = w.lower()
            if lw in ('revenue','rev','arr','users','customers','growth','market','marketshare','cost','sales','profit','loss','margin','revenue:'):
                metric = lw
                break
        if not metric:
            # fallback to kind + a truncated snippet
            metric = f.kind + ":" + (snippet[:30].strip().lower())
        groups[metric].append(f)
    return groups

# --------------------------- Comparison ---------------------------


def relative_difference(a: float, b: float) -> float:
    try:
        if a == 0 and b == 0:
            return 0.0
        return abs(a - b) / (abs(a) + abs(b))
    except Exception:
        return float('inf')


def rule_based_compare(fa: Fact, fb: Fact) -> Optional[Issue]:
    """Rule-based comparison for numeric and date facts. Returns an Issue if detected, else None."""
    # Only compare facts of same `kind` or compatible kinds
    if fa.kind in ('currency','number','percent') and fb.kind in ('currency','number','percent'):
        va = fa.normalized['value'] if isinstance(fa.normalized, dict) and 'value' in fa.normalized else fa.normalized
        vb = fb.normalized['value'] if isinstance(fb.normalized, dict) and 'value' in fb.normalized else fb.normalized
        if va is None or vb is None:
            return None
        try:
            va = float(va)
            vb = float(vb)
        except Exception:
            return None
        rel = relative_difference(va, vb)
        if rel > 0.05:  # threshold: >5% difference
            details = {'value_a': va, 'value_b': vb, 'relative_diff': rel}
            explanation = f"Numeric mismatch: {va} vs {vb} (relative diff {rel:.2%})"
            return Issue(type='numeric_mismatch', slide_a=fa.slide_no, slide_b=fb.slide_no,
                         snippet_a=fa.raw, snippet_b=fb.raw, details=details, confidence=0.6, explanation=explanation)
        else:
            return None

    # Date contradictions: if parsed years differ significantly
    if fa.kind == 'date' and fb.kind == 'date':
        try:
            y1 = int(fa.normalized[:4])
            y2 = int(fb.normalized[:4])
            if y1 != y2:
                details = {'date_a': fa.normalized, 'date_b': fb.normalized}
                explanation = f"Date mismatch: {fa.normalized} vs {fb.normalized}"
                return Issue(type='date_mismatch', slide_a=fa.slide_no, slide_b=fb.slide_no,
                             snippet_a=fa.raw, snippet_b=fb.raw, details=details, confidence=0.7, explanation=explanation)
        except Exception:
            return None

    # Claim-level heuristic contradictions (simple negation detection)
    if fa.kind == 'claim' and fb.kind == 'claim':
        a = fa.raw.lower()
        b = fb.raw.lower()
        # naive: look for simple antonyms
        contradictions = [('few', 'many'), ('low', 'high'), ('decrease', 'increase'), ('no competitors', 'competitive')]
        for p, q in contradictions:
            if p in a and q in b or q in a and p in b:
                details = {'claim_a': fa.raw, 'claim_b': fb.raw}
                explanation = f"Textual contradiction heuristic between slides {fa.slide_no} and {fb.slide_no}"
                return Issue(type='textual_contradiction', slide_a=fa.slide_no, slide_b=fb.slide_no,
                             snippet_a=fa.raw, snippet_b=fb.raw, details=details, confidence=0.5, explanation=explanation)

    return None

# --------------------------- LLM integration (Gemini) ---------------------------

GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
GEMINI_ENDPOINT = os.environ.get('GEMINI_ENDPOINT')  # user should provide the correct REST endpoint
def call_gemini_batch(pairs: list, model: str = 'gemini-2.5-flash', timeout: int = 20):
    """
    Send multiple fact-pair comparisons in one Gemini call.
    `pairs` is a list of tuples: [(FactA, FactB), ...]
    Returns: list of Issue objects or None for skipped ones.
    """
    if not GEMINI_API_KEY or not GEMINI_ENDPOINT:
        return []

    # Build a single prompt for all pairs
    prompt = (
        "You will compare multiple slide claim pairs and output ONLY a JSON array, "
        "one object per pair, in the order provided.\n"
        "Each object must have:\n"
        "{\n"
        "  \"decision\": \"consistent|contradictory|numeric_mismatch|uncertain\",\n"
        "  \"reason\": \"brief explanation\",\n"
        "  \"evidence\": [\"snippet from A\", \"snippet from B\"],\n"
        "  \"score\": 0.0\n"
        "}\n\n"
        "Pairs:\n"
    )
    for idx, (fa, fb) in enumerate(pairs):
        prompt += f"{idx+1}. A(slide {fa.slide_no}): {fa.raw}\n"
        prompt += f"   B(slide {fb.slide_no}): {fb.raw}\n"

    endpoint_url = f"{GEMINI_ENDPOINT}?key={GEMINI_API_KEY}"
    headers = {'Content-Type': 'application/json'}
    body = {
        "contents": [
            {"parts": [{"text": prompt}]}
        ]
    }

    try:
        r = requests.post(endpoint_url, headers=headers, json=body, timeout=timeout)
        if r.status_code >= 400:
            print(f"[WARN] Gemini API error {r.status_code}: {r.text}")
            return []

        data = r.json()
        text_output = ""
        if "candidates" in data and data["candidates"]:
            parts = data["candidates"][0]["content"].get("parts", [])
            text_parts = [p.get("text", "").strip() for p in parts if "text" in p]
            text_output = "\n".join([t for t in text_parts if t])

        if not text_output:
            print("[WARN] Gemini returned empty output for batch.")
            return []

        try:
            results = json.loads(text_output)
        except Exception as e:
            print(f"[WARN] Could not parse Gemini batch output: {e}")
            print("Raw model output:\n", text_output)
            return []

        issues = []
        for (fa, fb), res in zip(pairs, results):
            decision = res.get("decision")
            reason = res.get("reason")
            evidence = res.get("evidence", [])
            score = float(res.get("score", 0.0)) if res.get("score") is not None else 0.0
            if decision in ("numeric_mismatch", "contradictory"):
                itype = "numeric_mismatch" if decision == "numeric_mismatch" else "textual_contradiction"
                issues.append(Issue(
                    type=itype,
                    slide_a=fa.slide_no,
                    slide_b=fb.slide_no,
                    snippet_a=fa.raw,
                    snippet_b=fb.raw,
                    details={"llm_reason": reason, "evidence": evidence},
                    confidence=score,
                    explanation=reason
                ))
        return issues

    except Exception as e:
        print(f"[WARN] Gemini batch call failed: {e}")
        return []
# --------------------------- Main analyzer ---------------------------
class DeckAnalyzer:
    def __init__(self, gemini_enabled: bool = False):
        # Force-enable Gemini if env vars are present
        self.gemini_enabled = gemini_enabled and bool(GEMINI_API_KEY) and bool(GEMINI_ENDPOINT)
        print(f"[DEBUG] Gemini enabled: {self.gemini_enabled}, API key present: {bool(GEMINI_API_KEY)}")

    def analyze(self, pptx_path: str, run_ocr: bool = True) -> Dict[str, Any]:
        slides = extract_from_pptx(pptx_path)
        print(f"Extracted {len(slides)} slides.")

        # Run OCR on extracted images and append OCR text to slide raw_text
        if run_ocr:
            for s in slides:
                for img in s.image_paths:
                    ocr_txt = ocr_image(img)
                    if ocr_txt:
                        s.raw_text += "\n" + ocr_txt

        # Parse facts
        all_facts: List[Fact] = []
        for s in slides:
            facts = parse_facts_from_text(s.raw_text, s.slide_no)
            all_facts.extend(facts)
        print(f"Parsed {len(all_facts)} facts across slides.")

        # Group facts heuristically
        groups = group_facts_by_metric(all_facts)
        print(f"Formed {len(groups)} heuristic metric groups.")

        issues: List[Issue] = []

        # For each group, compare facts pairwise
        for metric, facts in groups.items():
            if len(facts) < 2:
                continue

            # First pass: rule-based check
            batch_pairs = []
            for i in range(len(facts)):
                for j in range(i + 1, len(facts)):
                    fa = facts[i]
                    fb = facts[j]

                    rb = rule_based_compare(fa, fb)
                    if rb:
                        issues.append(rb)
                    else:
                        if self.gemini_enabled:
                            batch_pairs.append((fa, fb))
                        else:
                            # fallback: simple textual heuristic if no Gemini
                            la = fa.raw.lower()
                            lb = fb.raw.lower()
                            contradictions = [
                                ('few', 'many'),
                                ('no competitors', 'competitive'),
                                ('not profitable', 'profitable')
                            ]
                            for p, q in contradictions:
                                if (p in la and q in lb) or (q in la and p in lb):
                                    issues.append(Issue(
                                        type='textual_contradiction',
                                        slide_a=fa.slide_no,
                                        slide_b=fb.slide_no,
                                        snippet_a=fa.raw,
                                        snippet_b=fb.raw,
                                        details={'heuristic': (p, q)},
                                        confidence=0.45,
                                        explanation=f"Heuristic contradiction: '{p}' vs '{q}'"
                                    ))

            # Second pass: send batched Gemini calls
            if self.gemini_enabled and batch_pairs:
                for k in range(0, len(batch_pairs), 5):  # 5 pairs per call
                    sub_batch = batch_pairs[k:k + 5]
                    issues.extend(call_gemini_batch(sub_batch))

        # Deduplicate issues
        uniq = []
        seen = set()
        for it in issues:
            key = (it.type, min(it.slide_a, it.slide_b), max(it.slide_a, it.slide_b), it.snippet_a, it.snippet_b)
            if key not in seen:
                seen.add(key)
                uniq.append(it)

        # Build report
        report = {
            'meta': {
                'analyzed_file': os.path.abspath(pptx_path),
                'slides': len(slides),
                'facts': len(all_facts),
                'groups': len(groups),
                'gemini_used': self.gemini_enabled,
            },
            'issues': [asdict(x) for x in uniq]
        }
        return report

# --------------------------- CLI ---------------------------

def main():
    parser = argparse.ArgumentParser(description='Analyze a PPTX for factual/logical inconsistencies.')
    parser.add_argument('--input', '-i', required=True, help='Path to .pptx file')
    parser.add_argument('--out', '-o', default='report.json', help='Path to JSON report output')
    parser.add_argument('--no-ocr', dest='ocr', action='store_false', help='Disable OCR on slide images')
    parser.add_argument('--use-gemini', dest='use_gemini', action='store_true', help='Enable Gemini comparator (requires GEMINI_API_KEY & GEMINI_ENDPOINT env vars)')

    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Input file {args.input} not found.")
        sys.exit(2)

    analyzer = DeckAnalyzer(gemini_enabled=args.use_gemini)
    report = analyzer.analyze(args.input, run_ocr=args.ocr)

    with open(args.out, 'w', encoding='utf-8') as f:
        json.dump(report, f, indent=2, ensure_ascii=False)

    print(f"Report written to {args.out}")
    # also print a brief human summary
    issues = report.get('issues', [])
    if not issues:
        print("No issues detected (rule-based checks). Consider enabling Gemini for semantic checks.")
    else:
        print(f"Detected {len(issues)} issues:")
        for it in issues:
            print(f"- [{it['type']}] slides {it['slide_a']} vs {it['slide_b']}: {it['explanation']} (confidence {it['confidence']})")

if __name__ == '__main__':
    main()
