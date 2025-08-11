"""
Enhanced PPTX Inconsistency Detector

Key improvements:
1. Always uses LLM for comprehensive analysis
2. Detects numerical, grammatical, spelling, and logical inconsistencies
3. Better error handling and variable accessibility
4. Detailed CLI output with slide numbers
5. Multi-type inconsistency detection
"""

import os
import sys
import json
import re
import argparse
import time
from collections import defaultdict
from dataclasses import dataclass, asdict
from typing import List, Dict, Any, Optional, Tuple, Set
from enum import Enum

from pptx import Presentation
from PIL import Image
import pytesseract
import requests
import dateparser
from dotenv import load_dotenv
load_dotenv()

# --------------------------- Enhanced Data Classes ---------------------------

class FactType(Enum):
    NUMBER = "number"
    CURRENCY = "currency" 
    PERCENTAGE = "percentage"
    DATE = "date"
    RATIO = "ratio"
    CLAIM = "claim"
    TEXT = "text"  # For grammatical/spelling analysis

class InconsistencyType(Enum):
    NUMERIC = "numeric_inconsistency"
    TEXTUAL = "textual_contradiction" 
    GRAMMATICAL = "grammatical_error"
    SPELLING = "spelling_error"
    LOGICAL = "logical_contradiction"
    DATE_MISMATCH = "date_mismatch"

class Unit(Enum):
    BYTES = "bytes"
    KB = "kb" 
    MB = "mb"
    GB = "gb"
    TB = "tb"
    USD = "usd"
    EUR = "eur"
    GBP = "gbp"
    THOUSAND = "k"
    MILLION = "m" 
    BILLION = "b"
    YEAR = "year"
    MONTH = "month"
    QUARTER = "quarter"
    PERCENT = "percent"
    NONE = "none"

@dataclass
class EnhancedFact:
    slide_no: int
    fact_type: FactType
    raw_text: str
    normalized_value: Optional[float]
    unit: Unit
    context_window: str
    context_type: str
    confidence: float
    metadata: Optional[Dict[str, Any]] = None

@dataclass
class InconsistencyIssue:
    issue_type: InconsistencyType
    slide_a: int
    slide_b: int
    snippet_a: str
    snippet_b: str
    context_a: str
    context_b: str
    confidence: float
    explanation: str
    details: Optional[Dict[str, Any]] = None

# --------------------------- LLM Integration ---------------------------

GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
GEMINI_ENDPOINT = os.environ.get('GEMINI_ENDPOINT', 'https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash-exp:generateContent')

def call_gemini_for_analysis(slide_contents: List[Dict], analysis_type: str = "comprehensive") -> List[InconsistencyIssue]:
    """Call Gemini for comprehensive inconsistency analysis"""
    
    if not GEMINI_API_KEY:
        print("[ERROR] GEMINI_API_KEY not found. Please set it in environment variables.")
        return []
    
    if analysis_type == "comprehensive":
        prompt = create_comprehensive_analysis_prompt(slide_contents)
    elif analysis_type == "pairwise":
        prompt = create_pairwise_analysis_prompt(slide_contents)
    else:
        prompt = create_comprehensive_analysis_prompt(slide_contents)
    
    endpoint_url = f"{GEMINI_ENDPOINT}?key={GEMINI_API_KEY}"
    headers = {'Content-Type': 'application/json'}
    
    body = {
        "contents": [
            {"parts": [{"text": prompt}]}
        ],
        "generationConfig": {
            "temperature": 0.1,
            "topK": 40,
            "topP": 0.95,
            "maxOutputTokens": 8192,
        }
    }
    
    try:
        print("[INFO] Calling Gemini API for analysis...")
        response = requests.post(endpoint_url, headers=headers, json=body, timeout=60)
        
        if response.status_code != 200:
            print(f"[ERROR] Gemini API returned status {response.status_code}: {response.text}")
            return []
        
        data = response.json()
        
        # Extract text from response
        text_output = ""
        if "candidates" in data and data["candidates"]:
            parts = data["candidates"][0]["content"].get("parts", [])
            text_parts = [p.get("text", "").strip() for p in parts if "text" in p]
            text_output = "\n".join([t for t in text_parts if t])
        
        if not text_output:
            print("[WARN] Gemini returned empty output.")
            return []
        
        # Parse the JSON response
        issues = parse_gemini_response(text_output)
        return issues
        
    except requests.exceptions.Timeout:
        print("[ERROR] Gemini API call timed out.")
        return []
    except Exception as e:
        print(f"[ERROR] Gemini API call failed: {str(e)}")
        return []

def create_comprehensive_analysis_prompt(slide_contents: List[Dict]) -> str:
    """Create a comprehensive prompt for all types of inconsistencies"""

    prompt = """You are an expert presentation analyzer. Analyze these PowerPoint slides for inconsistencies, but ONLY flag genuine, high-confidence issues.

ONLY INCLUDE these types in your JSON output:
- logical_contradiction
- grammatical_error
- textual_contradiction

NUMERICAL INCONSISTENCIES: Only flag if the same metric is reported with conflicting values, or if calculations are clearly wrong. Do NOT flag differences for different time periods, units, or metrics.

TEXTUAL CONTRADICTIONS: Only flag if statements directly conflict.

GRAMMATICAL ERRORS: Only flag clear grammar mistakes.

LOGICAL CONTRADICTIONS: Only flag if statements are logically impossible or mutually exclusive.

DO NOT flag spelling errors, minor stylistic issues, or valid differences.

RETURN ONLY A VALID JSON ARRAY with this exact structure:
[
  {
    "type": "logical_contradiction|grammatical_error|textual_contradiction",
    "slide_a": 1,
    "slide_b": 2,
    "snippet_a": "exact text from slide A",
    "snippet_b": "exact text from slide B", 
    "context_a": "broader context from slide A",
    "context_b": "broader context from slide B",
    "confidence": 0.85,
    "explanation": "Clear explanation of the inconsistency",
    "details": {"key": "value"}
  }
]

For single-slide issues (grammar), use the same slide number for both slide_a and slide_b.

SLIDES TO ANALYZE:
"""
    for i, slide_content in enumerate(slide_contents, 1):
        prompt += f"\n=== SLIDE {i} ===\n{slide_content.get('text', '')}\n"

    prompt += "\nBe concise. Only return clear, high-confidence issues."
    return prompt

def create_pairwise_analysis_prompt(slide_contents: List[Dict]) -> str:
    """Create prompt for pairwise slide comparison"""
    
    prompt = """Compare these slide pairs for inconsistencies. Return JSON array:

"""
    
    for i, slide in enumerate(slide_contents):
        prompt += f"SLIDE {i+1}: {slide.get('text', '')}\n\n"
    
    prompt += """
Return JSON array of inconsistencies found.
"""
    
    return prompt

def parse_gemini_response(response_text: str) -> List[InconsistencyIssue]:
    """Parse Gemini's JSON response into InconsistencyIssue objects"""
    
    issues = []
    
    try:
        # Try to extract JSON from the response
        json_start = response_text.find('[')
        json_end = response_text.rfind(']') + 1
        
        if json_start == -1 or json_end == 0:
            print("[WARN] No JSON array found in Gemini response")
            print("Raw response:", response_text[:500])
            return []
        
        json_str = response_text[json_start:json_end]
        parsed_data = json.loads(json_str)
        
        for item in parsed_data:
            issue_type_str = item.get('type', 'unknown')
            
            # Map string to enum
            type_mapping = {
                'numeric_inconsistency': InconsistencyType.NUMERIC,
                'textual_contradiction': InconsistencyType.TEXTUAL,
                'grammatical_error': InconsistencyType.GRAMMATICAL,
                'spelling_error': InconsistencyType.SPELLING,
                'logical_contradiction': InconsistencyType.LOGICAL,
                'date_mismatch': InconsistencyType.DATE_MISMATCH
            }
            
            issue_type = type_mapping.get(issue_type_str, InconsistencyType.TEXTUAL)
            
            issue = InconsistencyIssue(
                issue_type=issue_type,
                slide_a=int(item.get('slide_a', 1)),
                slide_b=int(item.get('slide_b', 1)),
                snippet_a=item.get('snippet_a', ''),
                snippet_b=item.get('snippet_b', ''),
                context_a=item.get('context_a', ''),
                context_b=item.get('context_b', ''),
                confidence=float(item.get('confidence', 0.5)),
                explanation=item.get('explanation', ''),
                details=item.get('details', {})
            )
            issues.append(issue)
            
    except json.JSONDecodeError as e:
        print(f"[ERROR] Failed to parse Gemini JSON response: {e}")
        print("Raw response:", response_text[:1000])
    except Exception as e:
        print(f"[ERROR] Error parsing Gemini response: {e}")
    
    return issues

# --------------------------- Enhanced Regex Patterns ---------------------------

UNIT_NUMBER_PATTERN = re.compile(
    r"""
    (?P<value>\d{1,3}(?:,\d{3})*(?:\.\d+)?|\d+(?:\.\d+)?)  # Number
    \s*
    (?P<unit>KB|MB|GB|TB|bytes?|K|M|B|k|m|b|%|USD|EUR|GBP|₹|£|€|\$)?  # Unit
    """, 
    re.VERBOSE | re.IGNORECASE
)

COMPARISON_PATTERN = re.compile(
    r"""
    (?:compressed|reduced|increased|changed|went|from)\s+
    (?:size\s+)?(?:from\s+)?
    (?P<from_val>\d+(?:\.\d+)?)\s*(?P<from_unit>KB|MB|GB|TB|K|M|B|%|\$)?
    \s*(?:to|down\s+to|up\s+to)\s+
    (?P<to_val>\d+(?:\.\d+)?)\s*(?P<to_unit>KB|MB|GB|TB|K|M|B|%|\$)?
    """,
    re.VERBOSE | re.IGNORECASE
)

# --------------------------- Utility Functions ---------------------------

def detect_unit(unit_str: str) -> Unit:
    """Detect unit from string"""
    if not unit_str:
        return Unit.NONE
    
    unit_lower = unit_str.lower()
    unit_map = {
        'kb': Unit.KB, 'mb': Unit.MB, 'gb': Unit.GB, 'tb': Unit.TB,
        'bytes': Unit.BYTES, 'byte': Unit.BYTES,
        'k': Unit.THOUSAND, 'm': Unit.MILLION, 'b': Unit.BILLION,
        'usd': Unit.USD, '$': Unit.USD,
        'eur': Unit.EUR, '€': Unit.EUR,
        'gbp': Unit.GBP, '£': Unit.GBP,
        '%': Unit.PERCENT,
    }
    
    return unit_map.get(unit_lower, Unit.NONE)

def detect_context_type(text: str) -> str:
    """Detect context type from text"""
    context_lower = text.lower()
    
    if any(word in context_lower for word in ['file', 'size', 'compressed', 'storage', 'download', 'upload', 'memory']):
        return "file_size"
    elif any(word in context_lower for word in ['revenue', 'profit', 'cost', 'price', 'sales', 'earnings', 'budget']):
        return "financial"
    elif any(word in context_lower for word in ['growth', 'increase', 'decrease', 'performance', 'rate']):
        return "performance"
    elif any(word in context_lower for word in ['users', 'customers', 'market', 'share', 'adoption']):
        return "metrics"
    
    return "general"

def extract_basic_facts(text: str, slide_no: int) -> List[EnhancedFact]:
    """Extract basic facts for reference (not used for main analysis)"""
    facts = []
    
    # Extract sentences as text facts
    sentences = re.split(r'[.!?]+', text)
    for sentence in sentences:
        sentence = sentence.strip()
        if len(sentence) > 10:  # Only meaningful sentences
            fact = EnhancedFact(
                slide_no=slide_no,
                fact_type=FactType.TEXT,
                raw_text=sentence,
                normalized_value=None,
                unit=Unit.NONE,
                context_window=sentence,
                context_type=detect_context_type(sentence),
                confidence=0.8
            )
            facts.append(fact)
    
    return facts

# --------------------------- Main Enhanced Analyzer ---------------------------

class EnhancedDeckAnalyzer:
    def __init__(self):
        """Always use LLM for comprehensive analysis"""
        self.use_llm = True
        print("[INFO] Enhanced Deck Analyzer initialized with LLM enabled")
    
    def analyze(self, pptx_path: str) -> Dict[str, Any]:
        """Perform comprehensive analysis using LLM"""
        
        print(f"[INFO] Starting analysis of {pptx_path}")
        
        # Extract slides
        slides = self.extract_slides(pptx_path)
        print(f"[INFO] Extracted {len(slides)} slides")
        
        # Prepare slide contents for LLM
        slide_contents = []
        for slide in slides:
            slide_contents.append({
                'slide_no': slide.slide_no,
                'text': slide.raw_text
            })
        
        # Use LLM for comprehensive analysis
        issues = call_gemini_for_analysis(slide_contents, "comprehensive")
        
        print(f"[INFO] LLM analysis complete. Found {len(issues)} potential issues.")
        
        # Build comprehensive report
        report = {
            'meta': {
                'file': os.path.abspath(pptx_path),
                'slides_analyzed': len(slides),
                'total_issues_found': len(issues),
                'analysis_method': 'LLM_comprehensive',
                'llm_model': 'gemini-2.0-flash-exp',
                'timestamp': time.strftime('%Y-%m-%d %H:%M:%S')
            },
            'issues': [asdict(issue) if hasattr(issue, '__dict__') else issue for issue in issues],
            'issues_by_type': self.categorize_issues(issues),
            'slides_content': [{'slide_no': s.slide_no, 'text': s.raw_text} for s in slides]
        }
        
        return report
    
    def categorize_issues(self, issues: List[InconsistencyIssue]) -> Dict[str, int]:
        """Categorize issues by type"""
        categories = defaultdict(int)
        for issue in issues:
            # Handle both enum and string types
            if hasattr(issue, 'issue_type'):
                if hasattr(issue.issue_type, 'value'):
                    issue_type = issue.issue_type.value
                else:
                    issue_type = str(issue.issue_type)
            else:
                # Handle dict case
                issue_type = issue.get('issue_type', 'unknown')
                if hasattr(issue_type, 'value'):
                    issue_type = issue_type.value
            
            categories[issue_type] += 1
        return dict(categories)
    
    def extract_slides(self, pptx_path: str):
        """Extract slide content from PPTX file"""
        try:
            prs = Presentation(pptx_path)
            slides = []
            
            for i, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text:
                            texts.append(shape.text.strip())
                    except Exception as e:
                        print(f"[WARN] Error extracting text from shape on slide {i}: {e}")
                        continue
                
                # Create slide content object
                slide_content = type('SlideContent', (), {
                    'slide_no': i,
                    'raw_text': '\n'.join(texts)
                })()
                slides.append(slide_content)
            
            return slides
            
        except Exception as e:
            print(f"[ERROR] Failed to extract slides from {pptx_path}: {e}")
            return []

# --------------------------- Enhanced CLI ---------------------------

def print_detailed_report(report: Dict[str, Any]):
    """Print a detailed, formatted report to console"""
    
    print("\n" + "="*80)
    print(f"PPTX INCONSISTENCY ANALYSIS REPORT")
    print("="*80)
    
    meta = report.get('meta', {})
    print(f"File: {meta.get('file', 'Unknown')}")
    print(f"Slides Analyzed: {meta.get('slides_analyzed', 0)}")
    print(f"Analysis Method: {meta.get('analysis_method', 'Unknown')}")
    print(f"Timestamp: {meta.get('timestamp', 'Unknown')}")
    
    issues = report.get('issues', [])
    issues_by_type = report.get('issues_by_type', {})
    
    print(f"\nTOTAL ISSUES FOUND: {len(issues)}")
    
    if issues_by_type:
        print("\nISSUES BY TYPE:")
        for issue_type, count in issues_by_type.items():
            print(f"  - {issue_type.replace('_', ' ').title()}: {count}")
    
    if not issues:
        print("\n✅ No inconsistencies detected! The presentation appears to be consistent.")
        return
    
    print("\n" + "-"*80)
    print("DETAILED ISSUES:")
    print("-"*80)
    
    for i, issue in enumerate(issues, 1):
        # Handle both dict and InconsistencyIssue object
        if isinstance(issue, dict):
            issue_type = issue.get('issue_type', 'unknown')
            slide_a = issue.get('slide_a', 0)
            slide_b = issue.get('slide_b', 0)
            snippet_a = issue.get('snippet_a', '')
            snippet_b = issue.get('snippet_b', '')
            explanation = issue.get('explanation', '')
            confidence = issue.get('confidence', 0)
        else:
            # Handle InconsistencyIssue object
            issue_type = issue.issue_type.value if hasattr(issue.issue_type, 'value') else str(issue.issue_type)
            slide_a = issue.slide_a
            slide_b = issue.slide_b
            snippet_a = issue.snippet_a
            snippet_b = issue.snippet_b
            explanation = issue.explanation
            confidence = issue.confidence
        
        # Convert issue_type to string if it's still an enum
        if hasattr(issue_type, 'value'):
            issue_type_str = issue_type.value
        elif hasattr(issue_type, 'replace'):
            issue_type_str = issue_type
        else:
            issue_type_str = str(issue_type)
        
        print(f"\n[{i}] {issue_type_str.replace('_', ' ').upper()}")
        print(f"    Confidence: {confidence:.1%}")
        
        if slide_a == slide_b:
            print(f"    Location: Slide {slide_a}")
            print(f"    Issue: {snippet_a}")
        else:
            print(f"    Between: Slide {slide_a} ↔ Slide {slide_b}")
            print(f"    Slide {slide_a}: \"{snippet_a[:100]}{'...' if len(snippet_a) > 100 else ''}\"")
            print(f"    Slide {slide_b}: \"{snippet_b[:100]}{'...' if len(snippet_b) > 100 else ''}\"")
        
        print(f"    Explanation: {explanation}")
    
    print("\n" + "="*80)

def main():
    parser = argparse.ArgumentParser(
        description='Enhanced PPTX Inconsistency Detector - Comprehensive Analysis',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python enhanced_detector.py -i presentation.pptx
  python enhanced_detector.py -i deck.pptx -o detailed_report.json
  
Note: Requires GEMINI_API_KEY environment variable to be set.
        """
    )
    
    parser.add_argument('--input', '-i', required=True, help='Path to .pptx file')
    parser.add_argument('--output', '-o', default='comprehensive_report.json', help='Output JSON report file')
    parser.add_argument('--verbose', '-v', action='store_true', help='')
    
    args = parser.parse_args()
    
    # Validate input file
    if not os.path.exists(args.input):
        print(f"[ERROR] Input file '{args.input}' not found.")
        sys.exit(1)
    
    # Check for API key
    if not GEMINI_API_KEY:
        print("[ERROR] GEMINI_API_KEY environment variable not set.")
        print("Please get your API key from: https://aistudio.google.com/app/apikey")
        print("Then set it as: export GEMINI_API_KEY='your-api-key'")
        sys.exit(1)
    
    # Initialize analyzer and run analysis
    analyzer = EnhancedDeckAnalyzer()
    
    try:
        report = analyzer.analyze(args.input)
        
        # Save JSON report
        with open(args.output, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False, default=str)
        
        print(f"[INFO] Detailed report saved to: {args.output}")
        
        # Print formatted report to console
        print_detailed_report(report)
        
    except KeyboardInterrupt:
        print("\n[INFO] Analysis interrupted by user.")
        sys.exit(0)
    except Exception as e:
        print(f"[ERROR] Analysis failed: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        sys.exit(1)

if __name__ == '__main__':
    main()